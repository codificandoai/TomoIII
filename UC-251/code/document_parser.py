"""Extracción de texto, tablas y metadatos de documentos empresariales.

Formatos soportados: PDF, DOCX, PPTX, HTML, imágenes (OCR), TXT/Markdown.
Cada extractor es opcional: si falta la librería se cae a lectura de texto plano
o a OCR stub, garantizando que el pipeline siga operativo.
"""
from __future__ import annotations

import logging
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional

from models import Document
from text_utils import clean_text, checksum_normalized

logger = logging.getLogger("uc251-parser")


class DocumentParser:
    """Parser multimodal con extracción de texto estructurado."""

    def __init__(self, ocr_enabled: bool = True):
        self.ocr_enabled = ocr_enabled

    def parse(
        self,
        path_or_text: str,
        doc_type: Optional[str] = None,
        title: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> Document:
        """Parsea un archivo o texto plano y devuelve un Document."""
        metadata = metadata or {}
        is_file = os.path.isfile(path_or_text)
        if is_file:
            p = Path(path_or_text)
            ext = (doc_type or p.suffix.lstrip(".")).lower()
            raw_bytes = p.read_bytes()
            detected_type = self._detect_mime_or_ext(p, raw_bytes, ext)
            content, extracted_tables, headings = self._extract_from_file(
                p, raw_bytes, detected_type
            )
            source = str(p.resolve())
            doc_id = checksum_normalized(f"{source}:{raw_bytes.hex()}")
        else:
            detected_type = (doc_type or "txt").lower()
            content, extracted_tables, headings = path_or_text, [], []
            source = metadata.get("source", "inline-text")
            doc_id = checksum_normalized(f"{source}:{path_or_text}")

        content = clean_text(content)
        if not content:
            logger.warning("Documento resultó vacío tras la limpieza: %s", source)

        full_metadata = {
            **metadata,
            "extracted_tables": extracted_tables,
            "extracted_headings": headings,
            "parser": detected_type,
        }
        return Document(
            doc_id=doc_id,
            source=source,
            title=title or source.split("/")[-1],
            content=content,
            doc_type=detected_type,
            metadata=full_metadata,
            checksum=checksum_normalized(content),
        )

    @staticmethod
    def _detect_mime_or_ext(path: Path, raw_bytes: bytes, ext: str) -> str:
        if raw_bytes.startswith(b"%PDF"):
            return "pdf"
        if raw_bytes[:4] == b"PK\x03\x04":
            # DOCX/PPTX/XLSX comparten firma ZIP
            name_lower = path.name.lower()
            if name_lower.endswith(".docx"):
                return "docx"
            if name_lower.endswith(".pptx"):
                return "pptx"
            return ext or "unknown"
        return ext or "txt"

    def _extract_from_file(
        self, path: Path, raw_bytes: bytes, doc_type: str
    ) -> tuple[str, List[Dict[str, Any]], List[str]]:
        """Devuelve (texto, tablas, encabezados)."""
        if doc_type == "pdf":
            return self._extract_pdf(path)
        if doc_type == "docx":
            return self._extract_docx(path)
        if doc_type == "pptx":
            return self._extract_pptx(path)
        if doc_type in ("html", "htm"):
            return self._extract_html(path)
        if doc_type in ("png", "jpg", "jpeg", "tiff", "bmp", "gif"):
            return self._extract_image(path)
        return self._extract_text(path)

    def _extract_pdf(self, path: Path) -> tuple[str, List[Dict[str, Any]], List[str]]:
        try:
            import pypdf  # noqa: F401
            from pypdf import PdfReader
        except Exception as exc:  # pragma: no cover - environment dependent
            logger.warning("pypdf no disponible: %s", exc)
            return self._extract_text(path)

        reader = PdfReader(str(path))
        parts: List[str] = []
        tables: List[Dict[str, Any]] = []
        headings: List[str] = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            parts.append(f"\n--- Página {i} ---\n{text}")
            # pypdf básico no extrae tablas; dejamos un placeholder estructurado
            tables.extend(
                [{"page": i, "type": "pdf_table", "caption": f"Tabla página {i}"}]
            )
            for h in re.findall(r"^(?:[A-Z][A-Z\s]{2,}|\d+\.\s+.+)$", text, re.MULTILINE):
                if h.strip() and h not in headings:
                    headings.append(h.strip())
        return clean_text("\n".join(parts)), tables, headings

    def _extract_docx(self, path: Path) -> tuple[str, List[Dict[str, Any]], List[str]]:
        try:
            from docx import Document as DocxDocument
        except Exception as exc:  # pragma: no cover
            logger.warning("python-docx no disponible: %s", exc)
            return self._extract_text(path)

        doc = DocxDocument(str(path))
        paragraphs: List[str] = []
        tables: List[Dict[str, Any]] = []
        headings = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            paragraphs.append(text)
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                headings.append(text)
        for i, table in enumerate(doc.tables, start=1):
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            if rows:
                tables.append({"index": i, "type": "docx_table", "rows": rows})
                paragraphs.append(f"[Tabla {i}]\n" + "\n".join(" | ".join(r) for r in rows))
        return clean_text("\n\n".join(paragraphs)), tables, headings

    def _extract_pptx(self, path: Path) -> tuple[str, List[Dict[str, Any]], List[str]]:
        try:
            from pptx import Presentation
        except Exception as exc:  # pragma: no cover
            logger.warning("python-pptx no disponible: %s", exc)
            return self._extract_text(path)

        prs = Presentation(str(path))
        slides: List[str] = []
        tables: List[Dict[str, Any]] = []
        headings: List[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_texts.append(shape.text)
                if shape.has_table:
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    tables.append({"slide": i, "type": "pptx_table", "rows": rows})
            text = "\n".join(t for t in slide_texts if t.strip())
            if text.strip():
                slides.append(f"--- Diapositiva {i} ---\n{text}")
                headings.extend([t for t in slide_texts if len(t) < 80 and t.strip().istitle()])
        return clean_text("\n\n".join(slides)), tables, headings

    def _extract_html(self, path: Path) -> tuple[str, List[Dict[str, Any]], List[str]]:
        try:
            from bs4 import BeautifulSoup
        except Exception as exc:  # pragma: no cover
            logger.warning("BeautifulSoup no disponible: %s", exc)
            return self._extract_text(path)

        html = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        headings = [h.get_text(strip=True) for h in soup.find_all(["h1", "h2", "h3", "h4"])]
        tables = []
        for i, table in enumerate(soup.find_all("table"), start=1):
            rows = []
            for tr in table.find_all("tr"):
                row = [td.get_text(" ", strip=True) for td in tr.find_all(["td", "th"])]
                if row:
                    rows.append(row)
            if rows:
                tables.append({"index": i, "type": "html_table", "rows": rows})
        text = soup.get_text("\n", strip=True)
        return clean_text(text), tables, headings

    def _extract_image(self, path: Path) -> tuple[str, List[Dict[str, Any]], List[str]]:
        if not self.ocr_enabled:
            return self._extract_text(path)
        try:
            from PIL import Image
            import pytesseract
        except Exception as exc:  # pragma: no cover
            logger.warning("OCR no disponible (PIL/pytesseract): %s", exc)
            return self._extract_text(path)

        try:
            img = Image.open(str(path))
            text = pytesseract.image_to_string(img)
        except Exception as exc:  # pragma: no cover
            logger.warning("Error OCR %s: %s", path, exc)
            text = ""
        return clean_text(text), [{"type": "ocr", "source": str(path)}], []

    def _extract_text(self, path: Path) -> tuple[str, List[Dict[str, Any]], List[str]]:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return clean_text(text), [], []
